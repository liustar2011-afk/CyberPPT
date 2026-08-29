"""Native, deterministic source extractors shared by script and strict profiles.

The extractors preserve source wording, ordering, locators, and stable IDs. They
do not summarize, rank, or interpret source content.
"""
from __future__ import annotations

from collections import defaultdict
from pathlib import Path
import re
from typing import Any

from .source_document_map import (
    _extract_binary,
    _extract_docx,
    _extract_text,
    _heading_id,
    _sha256_bytes,
    _stable_id,
    _unit,
)

_TEXT_SUFFIXES = frozenset({".txt", ".md", ".json", ".csv", ".tsv", ".yaml", ".yml"})
_MARKITDOWN_SUFFIXES = frozenset({".pdf", ".html", ".htm", ".rtf"})
_MARKDOWN_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")


def _clean_text(value: object) -> str:
    return "\n".join(line.rstrip() for line in str(value or "").splitlines()).strip()


def _extract_markitdown_fallback(
    path: Path,
    *,
    source_id: str,
    source_path: str,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, str]]]:
    try:
        from markitdown import MarkItDown
    except ImportError:
        units, headings, warnings = _extract_binary(
            path, source_id=source_id, source_path=source_path
        )
        warnings.append(
            {
                "code": "SOURCE_FORMAT_REQUIRES_MARKITDOWN",
                "message": f"{source_path} requires the matching optional MarkItDown format extra.",
            }
        )
        return units, headings, warnings
    try:
        converted = _clean_text(MarkItDown().convert(str(path)).text_content)
    except Exception as exc:  # optional converter families expose different errors
        units, headings, warnings = _extract_binary(
            path, source_id=source_id, source_path=source_path
        )
        warnings.append(
            {
                "code": "SOURCE_FORMAT_FALLBACK_FAILED",
                "message": f"{source_path}: {exc}",
            }
        )
        return units, headings, warnings
    if not converted:
        return [], [], [
            {
                "code": "SOURCE_FORMAT_FALLBACK_EMPTY",
                "message": f"{source_path} produced no text through MarkItDown.",
            }
        ]

    units: list[dict[str, Any]] = []
    headings: list[dict[str, Any]] = []
    duplicate_counts: dict[tuple[str, str], int] = defaultdict(int)
    heading_stack: list[dict[str, Any]] = []
    source_order = 0
    for line_number, raw in enumerate(converted.splitlines(), start=1):
        text = raw.strip()
        if not text:
            continue
        source_order += 1
        match = _MARKDOWN_HEADING_RE.match(text)
        if match:
            level = len(match.group(1))
            title = match.group(2).strip()
            while heading_stack and int(heading_stack[-1]["level"]) >= level:
                heading_stack.pop()
            parent_id = str(heading_stack[-1]["heading_id"]) if heading_stack else ""
            heading_id = _heading_id(source_id, level, title, parent_id, duplicate_counts)
            heading_path = [str(item["title"]) for item in heading_stack] + [title]
            unit_id = _stable_id(
                source_id,
                "heading",
                f"fallback:{line_number}\0{level}\0{title}",
                duplicate_counts,
            )
            units.append(
                _unit(
                    unit_id=unit_id,
                    source_id=source_id,
                    source_path=source_path,
                    kind="heading",
                    source_order=source_order,
                    text=title,
                    heading_id=heading_id,
                    heading_path=heading_path,
                    locator={"converted_line": line_number},
                    outline_level=level,
                )
            )
            heading = {
                "heading_id": heading_id,
                "source_id": source_id,
                "source_path": source_path,
                "title": title,
                "level": level,
                "parent_heading_id": parent_id or None,
                "source_order": source_order,
                "unit_id": unit_id,
                "heading_path": heading_path,
                "extraction_engine": "markitdown",
            }
            headings.append(heading)
            heading_stack.append(heading)
            continue
        heading_id = str(heading_stack[-1]["heading_id"]) if heading_stack else None
        heading_path = [str(item["title"]) for item in heading_stack]
        unit_id = _stable_id(
            source_id,
            "paragraph",
            f"fallback:{line_number}\0{text}",
            duplicate_counts,
        )
        units.append(
            _unit(
                unit_id=unit_id,
                source_id=source_id,
                source_path=source_path,
                kind="paragraph",
                source_order=source_order,
                text=text,
                heading_id=heading_id,
                heading_path=heading_path,
                locator={"converted_line": line_number},
                metadata={"extraction_engine": "markitdown"},
            )
        )
    return units, headings, [
        {
            "code": "SOURCE_FORMAT_EXTRACTED_WITH_MARKITDOWN",
            "message": f"{source_path} used the optional MarkItDown format fallback.",
        }
    ]


def extract_pptx_units(
    path: Path,
    *,
    source_id: str,
    source_path: str,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, str]]]:
    """Extract slide titles, text, tables, notes, and image placeholders."""

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(path)
    units: list[dict[str, Any]] = []
    headings: list[dict[str, Any]] = []
    warnings: list[dict[str, str]] = []
    duplicate_counts: dict[tuple[str, str], int] = defaultdict(int)
    source_order = 0
    image_count = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = _clean_text(title_shape.text if title_shape is not None else "")
        heading_id: str | None = None
        heading_path: list[str] = []
        if title:
            source_order += 1
            heading_id = _heading_id(source_id, 1, title, "", duplicate_counts)
            heading_path = [title]
            unit_id = _stable_id(
                source_id,
                "heading",
                f"slide:{slide_number}\0{title}",
                duplicate_counts,
            )
            units.append(
                _unit(
                    unit_id=unit_id,
                    source_id=source_id,
                    source_path=source_path,
                    kind="heading",
                    source_order=source_order,
                    text=title,
                    heading_id=heading_id,
                    heading_path=heading_path,
                    locator={"slide": slide_number, "shape": title_shape.shape_id},
                    outline_level=1,
                )
            )
            headings.append(
                {
                    "heading_id": heading_id,
                    "source_id": source_id,
                    "source_path": source_path,
                    "title": title,
                    "level": 1,
                    "parent_heading_id": None,
                    "source_order": source_order,
                    "unit_id": unit_id,
                    "heading_path": heading_path,
                }
            )

        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if getattr(shape, "has_chart", False):
                source_order += 1
                text = f"[chart on slide {slide_number}: {shape.name}; interpretation pending]"
                unit_id = _stable_id(
                    source_id,
                    "chart",
                    f"slide:{slide_number}\0shape:{shape.shape_id}\0{shape.name}",
                    duplicate_counts,
                )
                units.append(
                    _unit(
                        unit_id=unit_id,
                        source_id=source_id,
                        source_path=source_path,
                        kind="chart",
                        source_order=source_order,
                        text=text,
                        heading_id=heading_id,
                        heading_path=heading_path,
                        locator={"slide": slide_number, "shape": shape.shape_id},
                        metadata={"requires_visual_interpretation": True},
                    )
                )
                continue
            if getattr(shape, "has_table", False):
                for row_number, row in enumerate(shape.table.rows, start=1):
                    text = " | ".join(_clean_text(cell.text) for cell in row.cells).strip(" |")
                    if not text:
                        continue
                    source_order += 1
                    unit_id = _stable_id(
                        source_id,
                        "table_row",
                        f"slide:{slide_number}\0shape:{shape.shape_id}\0row:{row_number}\0{text}",
                        duplicate_counts,
                    )
                    units.append(
                        _unit(
                            unit_id=unit_id,
                            source_id=source_id,
                            source_path=source_path,
                            kind="table_row",
                            source_order=source_order,
                            text=text,
                            heading_id=heading_id,
                            heading_path=heading_path,
                            locator={"slide": slide_number, "shape": shape.shape_id, "row": row_number},
                        )
                    )
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                source_order += 1
                image_count += 1
                raw = shape.image.blob
                text = f"[image on slide {slide_number}: {shape.name}; visual interpretation pending]"
                unit_id = _stable_id(
                    source_id,
                    "image",
                    f"slide:{slide_number}\0shape:{shape.shape_id}\0{_sha256_bytes(raw)}",
                    duplicate_counts,
                )
                units.append(
                    _unit(
                        unit_id=unit_id,
                        source_id=source_id,
                        source_path=source_path,
                        kind="image",
                        source_order=source_order,
                        text=text,
                        heading_id=heading_id,
                        heading_path=heading_path,
                        locator={"slide": slide_number, "shape": shape.shape_id},
                        raw_sha256=_sha256_bytes(raw),
                        metadata={"requires_visual_interpretation": True},
                    )
                )
                continue
            text = _clean_text(shape.text if getattr(shape, "has_text_frame", False) else "")
            if text:
                source_order += 1
                unit_id = _stable_id(
                    source_id,
                    "paragraph",
                    f"slide:{slide_number}\0shape:{shape.shape_id}\0{text}",
                    duplicate_counts,
                )
                units.append(
                    _unit(
                        unit_id=unit_id,
                        source_id=source_id,
                        source_path=source_path,
                        kind="paragraph",
                        source_order=source_order,
                        text=text,
                        heading_id=heading_id,
                        heading_path=heading_path,
                        locator={"slide": slide_number, "shape": shape.shape_id},
                    )
                )

        notes = _clean_text(slide.notes_slide.notes_text_frame.text)
        if notes:
            source_order += 1
            unit_id = _stable_id(
                source_id,
                "speaker_note",
                f"slide:{slide_number}\0notes\0{notes}",
                duplicate_counts,
            )
            units.append(
                _unit(
                    unit_id=unit_id,
                    source_id=source_id,
                    source_path=source_path,
                    kind="speaker_note",
                    source_order=source_order,
                    text=notes,
                    heading_id=heading_id,
                    heading_path=heading_path,
                    locator={"slide": slide_number, "part": "speaker_notes"},
                )
            )

    if image_count:
        warnings.append(
            {
                "code": "SOURCE_IMAGE_SEMANTICS_PENDING",
                "message": f"{source_path} contains {image_count} image(s) requiring visual interpretation.",
            }
        )
    return units, headings, warnings


def extract_xlsx_units(
    path: Path,
    *,
    source_id: str,
    source_path: str,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, str]]]:
    """Extract worksheets and non-empty rows when optional openpyxl is present."""

    try:
        import openpyxl
    except ImportError:
        units, headings, warnings = _extract_binary(
            path, source_id=source_id, source_path=source_path
        )
        warnings.append(
            {
                "code": "SOURCE_XLSX_REQUIRES_OPENPYXL",
                "message": f"{source_path} requires optional openpyxl for native row extraction.",
            }
        )
        return units, headings, warnings

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=False)
    units: list[dict[str, Any]] = []
    headings: list[dict[str, Any]] = []
    duplicate_counts: dict[tuple[str, str], int] = defaultdict(int)
    source_order = 0
    try:
        for sheet_number, worksheet in enumerate(workbook.worksheets, start=1):
            title = str(worksheet.title)
            source_order += 1
            heading_id = _heading_id(source_id, 1, title, "", duplicate_counts)
            unit_id = _stable_id(
                source_id,
                "heading",
                f"sheet:{sheet_number}\0{title}",
                duplicate_counts,
            )
            units.append(
                _unit(
                    unit_id=unit_id,
                    source_id=source_id,
                    source_path=source_path,
                    kind="heading",
                    source_order=source_order,
                    text=title,
                    heading_id=heading_id,
                    heading_path=[title],
                    locator={"sheet": title},
                    outline_level=1,
                )
            )
            headings.append(
                {
                    "heading_id": heading_id,
                    "source_id": source_id,
                    "source_path": source_path,
                    "title": title,
                    "level": 1,
                    "parent_heading_id": None,
                    "source_order": source_order,
                    "unit_id": unit_id,
                    "heading_path": [title],
                }
            )
            for row_number, row in enumerate(worksheet.iter_rows(), start=1):
                values = [_clean_text(cell.value) for cell in row]
                if not any(values):
                    continue
                text = " | ".join(values).rstrip(" |")
                source_order += 1
                row_unit_id = _stable_id(
                    source_id,
                    "table_row",
                    f"sheet:{title}\0row:{row_number}\0{text}",
                    duplicate_counts,
                )
                units.append(
                    _unit(
                        unit_id=row_unit_id,
                        source_id=source_id,
                        source_path=source_path,
                        kind="table_row",
                        source_order=source_order,
                        text=text,
                        heading_id=heading_id,
                        heading_path=[title],
                        locator={"sheet": title, "row": row_number},
                    )
                )
                for cell in row:
                    if cell.data_type != "f" and not str(cell.value or "").startswith("="):
                        continue
                    formula = str(cell.value or "")
                    source_order += 1
                    formula_unit_id = _stable_id(
                        source_id,
                        "formula",
                        f"sheet:{title}\0cell:{cell.coordinate}\0{formula}",
                        duplicate_counts,
                    )
                    units.append(
                        _unit(
                            unit_id=formula_unit_id,
                            source_id=source_id,
                            source_path=source_path,
                            kind="formula",
                            source_order=source_order,
                            text=formula,
                            heading_id=heading_id,
                            heading_path=[title],
                            locator={"sheet": title, "cell": cell.coordinate},
                            metadata={"format": "xlsx_formula"},
                        )
                    )
    finally:
        workbook.close()
    return units, headings, []


def extract_source_file(
    path: Path,
    *,
    source_id: str,
    source_path: str,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, str]]]:
    """Dispatch one local source through the smallest native extractor."""

    suffix = path.suffix.casefold()
    if suffix == ".docx":
        return _extract_docx(path, source_id=source_id, source_path=source_path)
    if suffix == ".pptx":
        return extract_pptx_units(path, source_id=source_id, source_path=source_path)
    if suffix == ".xlsx":
        return extract_xlsx_units(path, source_id=source_id, source_path=source_path)
    if suffix in _TEXT_SUFFIXES:
        return _extract_text(path, source_id=source_id, source_path=source_path)
    if suffix in _MARKITDOWN_SUFFIXES:
        return _extract_markitdown_fallback(
            path, source_id=source_id, source_path=source_path
        )
    return _extract_binary(path, source_id=source_id, source_path=source_path)


__all__ = ["extract_pptx_units", "extract_source_file", "extract_xlsx_units"]
