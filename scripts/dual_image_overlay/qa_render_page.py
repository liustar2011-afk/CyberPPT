#!/usr/bin/env python3
"""One-shot render + structural QA for a single exported dual-image PPTX page.

Replaces the manual, multi-step "soffice -> pdftoppm -> crop -> eyeball" loop
with a single command that renders the page AND independently re-parses the
exported PPTX's actual shape geometry to check for overlaps and out-of-bounds
boxes. Screenshots are still produced for genuinely visual checks (contrast,
legibility), but geometry defects (the class of bug that actually took
multiple render-and-eyeball cycles to catch by hand) are now reported as a
structured pass/fail list before you look at a single pixel.

Usage:
    python3 scripts/dual_image_overlay/qa_render_page.py path/to/deck.pptx --out-dir qa_render
"""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu

SCRIPTS_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPTS_DIR.parents[1]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from scripts.dual_image_overlay.workspace_layout_qa import (  # noqa: E402
    check_page_layout_overlaps,
)


def _emu_to_px(value: int, *, dpi: int = 96) -> float:
    return round(float(value) / 914400.0 * dpi, 3)


def _shape_assignments(slide: Any, *, dpi: int) -> list[dict[str, Any]]:
    assignments: list[dict[str, Any]] = []
    for index, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
            continue
        assignments.append(
            {
                "text_index": index,
                "text": text,
                "final_bbox": {
                    "x": _emu_to_px(shape.left, dpi=dpi),
                    "y": _emu_to_px(shape.top, dpi=dpi),
                    "w": _emu_to_px(shape.width, dpi=dpi),
                    "h": _emu_to_px(shape.height, dpi=dpi),
                },
            }
        )
    return assignments


def _out_of_bounds_issues(
    assignments: list[dict[str, Any]], *, slide_width_px: float, slide_height_px: float, tolerance_px: float = 1.0
) -> list[dict[str, Any]]:
    issues = []
    for item in assignments:
        bbox = item["final_bbox"]
        right = bbox["x"] + bbox["w"]
        bottom = bbox["y"] + bbox["h"]
        if bbox["x"] < -tolerance_px or bbox["y"] < -tolerance_px or right > slide_width_px + tolerance_px or bottom > slide_height_px + tolerance_px:
            issues.append(
                {
                    "text_index": item["text_index"],
                    "text": item["text"],
                    "bbox": bbox,
                    "slide_size": {"width": slide_width_px, "height": slide_height_px},
                }
            )
    return issues


def check_pptx_geometry(pptx_path: Path, *, dpi: int = 96) -> dict[str, Any]:
    """Re-parse the actual exported PPTX shapes and check overlaps + bounds.

    This is deliberately independent of whatever `workspace_assignment.json`
    said upstream: it catches defects introduced anywhere in the chain,
    including a manual post-export patch that moved or resized a shape.
    """
    presentation = Presentation(str(pptx_path))
    slide_width_px = _emu_to_px(presentation.slide_width, dpi=dpi)
    slide_height_px = _emu_to_px(presentation.slide_height, dpi=dpi)
    slides_report: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        assignments = _shape_assignments(slide, dpi=dpi)
        overlap_report = check_page_layout_overlaps(assignments)
        out_of_bounds = _out_of_bounds_issues(
            assignments, slide_width_px=slide_width_px, slide_height_px=slide_height_px
        )
        slides_report.append(
            {
                "slide_index": slide_index,
                "text_box_count": len(assignments),
                "overlap_count": overlap_report["overlap_count"],
                "overlaps": overlap_report["overlaps"],
                "out_of_bounds_count": len(out_of_bounds),
                "out_of_bounds": out_of_bounds,
                "valid": overlap_report["valid"] and not out_of_bounds,
            }
        )
    return {
        "schema": "cyberppt.dual_image.qa_render_geometry.v1",
        "pptx": str(pptx_path),
        "slide_width_in": round(presentation.slide_width / 914400.0, 4),
        "slide_height_in": round(presentation.slide_height / 914400.0, 4),
        "slide_count": len(slides_report),
        "valid": all(slide["valid"] for slide in slides_report),
        "slides": slides_report,
    }


def render_to_png(pptx_path: Path, out_dir: Path, *, dpi: int = 150) -> list[Path]:
    """Render the pptx to PDF (LibreOffice) then to PNG (poppler). Best-effort."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        print(
            "Warning: soffice/pdftoppm not found on PATH; skipping render, geometry check still ran.",
            file=sys.stderr,
        )
        return []
    out_dir.mkdir(parents=True, exist_ok=True)
    render_env = os.environ.copy()
    if not render_env.get("FONTCONFIG_FILE"):
        for candidate in (
            Path("/opt/homebrew/etc/fonts/fonts.conf"),
            Path("/usr/local/etc/fonts/fonts.conf"),
            Path("/etc/fonts/fonts.conf"),
        ):
            if candidate.is_file():
                render_env["FONTCONFIG_FILE"] = str(candidate)
                break
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
        check=True,
        capture_output=True,
        env=render_env,
    )
    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    subprocess.run(
        [pdftoppm, "-jpeg", "-r", str(dpi), str(pdf_path), str(out_dir / "slide")],
        check=True,
        capture_output=True,
    )
    return sorted(out_dir.glob("slide*.jpg"))


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out-dir", type=Path, default=None, help="Directory for render + report output (default: alongside the pptx, in qa_render/)")
    parser.add_argument("--dpi", type=int, default=150)
    parser.add_argument("--no-render", action="store_true", help="Only run the geometry check; skip PDF/PNG rendering.")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    pptx_path = args.pptx.expanduser().resolve()
    if not pptx_path.is_file():
        print(f"Error: pptx not found: {pptx_path}", file=sys.stderr)
        return 1
    out_dir = (args.out_dir or pptx_path.parent / "qa_render").expanduser().resolve()

    geometry_report = check_pptx_geometry(pptx_path)
    out_dir.mkdir(parents=True, exist_ok=True)
    report_path = out_dir / f"{pptx_path.stem}_qa_geometry.json"
    report_path.write_text(json.dumps(geometry_report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    print(f"Slide size: {geometry_report['slide_width_in']}in x {geometry_report['slide_height_in']}in")
    for slide in geometry_report["slides"]:
        status = "OK" if slide["valid"] else "ISSUES"
        print(f"Slide {slide['slide_index']}: {slide['text_box_count']} text boxes - {status}")
        for overlap in slide["overlaps"]:
            print(
                f"  OVERLAP: [{overlap['text_index_a']}] {overlap['text_a']!r} <-> "
                f"[{overlap['text_index_b']}] {overlap['text_b']!r} "
                f"(area={overlap['overlap_area']})"
            )
        for oob in slide["out_of_bounds"]:
            print(f"  OUT OF BOUNDS: [{oob['text_index']}] {oob['text']!r} bbox={oob['bbox']}")
    print(f"Geometry report: {report_path}")

    if not args.no_render:
        images = render_to_png(pptx_path, out_dir, dpi=args.dpi)
        for image in images:
            print(f"Rendered: {image}")

    return 0 if geometry_report["valid"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
