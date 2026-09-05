from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from scripts.presentation_qa.render_page import check_pptx_geometry


def test_grouped_text_boxes_are_counted_recursively(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    first = group.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
    first.text_frame.text = "组内文本一"
    second = group.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(0.5))
    second.text_frame.text = "组内文本二"
    output = tmp_path / "grouped.pptx"
    presentation.save(output)

    report = check_pptx_geometry(output)

    assert report["valid"] is True
    assert report["slides"][0]["text_box_count"] == 2
