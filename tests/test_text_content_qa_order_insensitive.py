from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation
from pptx.util import Inches

from scripts.presentation_qa.text_content import build_text_content_qa


def _make_pptx(path: Path, texts: list[str]) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for index, text in enumerate(texts):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5 + index), Inches(2.0), Inches(0.4))
        box.text_frame.text = text
    presentation.save(str(path))


def test_order_insensitive_accepts_reordered_but_complete_text() -> None:
    with TemporaryDirectory() as directory:
        pptx_path = Path(directory) / "deck.pptx"
        _make_pptx(pptx_path, ["乙", "甲", "丙"])

        report = build_text_content_qa(pptx_path, ["甲", "乙", "丙"], order_sensitive=False)

        assert report["valid"] is True
        assert report["mismatches"] == []


def test_order_insensitive_still_flags_missing_text() -> None:
    with TemporaryDirectory() as directory:
        pptx_path = Path(directory) / "deck.pptx"
        _make_pptx(pptx_path, ["甲", "乙"])

        report = build_text_content_qa(pptx_path, ["甲", "乙", "丙"], order_sensitive=False)

        assert report["valid"] is False
        codes = [m["code"] for m in report["mismatches"]]
        assert "expected_text_missing_from_pptx" in codes


def test_order_insensitive_still_flags_unexpected_text() -> None:
    with TemporaryDirectory() as directory:
        pptx_path = Path(directory) / "deck.pptx"
        _make_pptx(pptx_path, ["甲", "乙", "丙"])

        report = build_text_content_qa(pptx_path, ["甲", "乙"], order_sensitive=False)

        assert report["valid"] is False
        codes = [m["code"] for m in report["mismatches"]]
        assert "unexpected_text_in_pptx" in codes


def test_order_sensitive_default_still_flags_reordering() -> None:
    with TemporaryDirectory() as directory:
        pptx_path = Path(directory) / "deck.pptx"
        _make_pptx(pptx_path, ["乙", "甲"])

        report = build_text_content_qa(pptx_path, ["甲", "乙"])

        assert report["valid"] is False
