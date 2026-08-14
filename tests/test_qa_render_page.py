from __future__ import annotations

import subprocess
from pathlib import Path
from tempfile import TemporaryDirectory

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

from scripts.dual_image_overlay import qa_render_page
from scripts.dual_image_overlay.qa_render_page import check_pptx_geometry


def _add_textbox(slide: object, text: str, *, left_in: float, top_in: float, w_in: float, h_in: float) -> None:
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in))  # type: ignore[attr-defined]
    box.text_frame.text = text


def test_clean_deck_reports_valid() -> None:
    with TemporaryDirectory() as directory:
        pptx_path = Path(directory) / "deck.pptx"
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _add_textbox(slide, "标题", left_in=0.5, top_in=0.3, w_in=3.0, h_in=0.5)
        _add_textbox(slide, "正文", left_in=0.5, top_in=1.0, w_in=3.0, h_in=0.5)
        presentation.save(str(pptx_path))

        report = check_pptx_geometry(pptx_path)

        assert report["valid"] is True
        assert round(report["slide_width_in"], 3) == 13.333
        assert report["slides"][0]["overlap_count"] == 0
        assert report["slides"][0]["out_of_bounds_count"] == 0


def test_detects_overlapping_shapes() -> None:
    with TemporaryDirectory() as directory:
        pptx_path = Path(directory) / "deck.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _add_textbox(slide, "授权管理", left_in=1.0, top_in=1.0, w_in=2.0, h_in=0.4)
        _add_textbox(slide, "访问主体与", left_in=1.05, top_in=1.1, w_in=2.0, h_in=0.4)
        presentation.save(str(pptx_path))

        report = check_pptx_geometry(pptx_path)

        assert report["valid"] is False
        assert report["slides"][0]["overlap_count"] == 1


def test_detects_out_of_bounds_shape() -> None:
    with TemporaryDirectory() as directory:
        pptx_path = Path(directory) / "deck.pptx"
        presentation = Presentation()
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(5.625)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(9.5), Inches(5.0), Inches(2.0), Inches(1.0))
        box.text_frame.text = "溢出文字"
        presentation.save(str(pptx_path))

        report = check_pptx_geometry(pptx_path)

        assert report["valid"] is False
        assert report["slides"][0]["out_of_bounds_count"] == 1


def test_ignores_empty_text_shapes() -> None:
    with TemporaryDirectory() as directory:
        pptx_path = Path(directory) / "deck.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(2.0), Inches(0.4))
        box.text_frame.text = ""
        presentation.save(str(pptx_path))

        report = check_pptx_geometry(pptx_path)

        assert report["slides"][0]["text_box_count"] == 0
        assert report["valid"] is True


def test_render_retries_bundled_office_after_homebrew_abort(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    pptx_path = tmp_path / "deck.pptx"
    pptx_path.write_bytes(b"pptx")
    output = tmp_path / "render"
    homebrew = Path("/opt/homebrew/bin/soffice")
    bundled = Path("/runtime/override/soffice")
    calls: list[list[str]] = []

    monkeypatch.setattr(qa_render_page, "_office_candidates", lambda: [homebrew, bundled])
    monkeypatch.setattr(
        qa_render_page.shutil,
        "which",
        lambda command: "/usr/bin/pdftoppm" if command == "pdftoppm" else None,
    )

    def run(command: list[str], **_kwargs: object) -> subprocess.CompletedProcess[str]:
        calls.append(command)
        if command[0] == str(homebrew):
            raise subprocess.CalledProcessError(134, command, output="homebrew stdout", stderr="homebrew aborted")
        if command[0] == str(bundled):
            (output / "deck.pdf").write_bytes(b"pdf")
        else:
            (output / "slide-1.jpg").write_bytes(b"jpg")
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(qa_render_page.subprocess, "run", run)

    assert qa_render_page.render_to_png(pptx_path, output) == [output / "slide-1.jpg"]
    assert [call[0] for call in calls] == [str(homebrew), str(bundled), "/usr/bin/pdftoppm"]
