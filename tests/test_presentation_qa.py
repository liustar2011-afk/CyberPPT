from __future__ import annotations

import unittest
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from scripts.presentation_qa.render_page import _render_with_officecli, check_pptx_geometry, render_to_png
from scripts.presentation_qa.text_content import build_text_content_qa
from cyberppt.officecli import OFFICECLI_VERSION, repository_officecli_path, resolve_officecli


class PresentationQaTests(unittest.TestCase):
    def _write_text_deck(self, path: Path) -> None:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        first = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
        first.text_frame.text = "生产运行需求："
        second = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(5), Inches(0.5))
        second.text_frame.text = "新能源接入后，安全保供需要及时数据支撑"
        presentation.save(str(path))

    def test_clean_deck_has_no_geometry_failures(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            presentation = Presentation()
            presentation.slide_width = Inches(13.333)
            presentation.slide_height = Inches(7.5)
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            first = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(3), Inches(0.5))
            first.text_frame.text = "标题"
            second = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(3), Inches(0.5))
            second.text_frame.text = "正文"
            presentation.save(str(pptx_path))

            report = check_pptx_geometry(pptx_path)

        self.assertTrue(report["valid"])
        self.assertEqual(0, report["slides"][0]["overlap_count"])
        self.assertEqual(0, report["slides"][0]["out_of_bounds_count"])

    def test_overlapping_text_boxes_fail_geometry_gate(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.4))
            first.text_frame.text = "授权管理"
            second = slide.shapes.add_textbox(Inches(1.05), Inches(1.1), Inches(2), Inches(0.4))
            second.text_frame.text = "访问主体"
            presentation.save(str(pptx_path))

            report = check_pptx_geometry(pptx_path)

        self.assertFalse(report["valid"])
        self.assertEqual(1, report["slides"][0]["overlap_count"])

    def test_default_renderer_prefers_officecli(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            out_dir = Path(directory) / "renders"
            with patch("scripts.presentation_qa.render_page._render_with_officecli", return_value=[out_dir / "slide-1.png"]) as officecli:
                with patch("scripts.presentation_qa.render_page._render_with_soffice") as soffice:
                    result = render_to_png(pptx_path, out_dir)

        self.assertEqual([out_dir / "slide-1.png"], result)
        officecli.assert_called_once_with(pptx_path, out_dir, dpi=150)
        soffice.assert_not_called()

    def test_repository_officecli_is_preferred_over_path(self) -> None:
        with patch("cyberppt.officecli.repository_officecli_path", return_value=Path("/repo/.tools/officecli")):
            with patch("cyberppt.officecli.Path.is_file", return_value=True):
                with patch("cyberppt.officecli.shutil.which") as which:
                    resolved = resolve_officecli()

        self.assertEqual(Path("/repo/.tools/officecli"), resolved)
        which.assert_not_called()

    def test_officecli_version_is_pinned(self) -> None:
        self.assertEqual("1.0.145", OFFICECLI_VERSION)
        self.assertIn("v1.0.145", str(repository_officecli_path()))

    def test_explicit_soffice_renderer_skips_officecli(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            out_dir = Path(directory) / "renders"
            with patch("scripts.presentation_qa.render_page._render_with_officecli") as officecli:
                with patch("scripts.presentation_qa.render_page._render_with_soffice", return_value=[out_dir / "slide-1.jpg"]) as soffice:
                    result = render_to_png(pptx_path, out_dir, renderer="soffice")

        self.assertEqual([out_dir / "slide-1.jpg"], result)
        officecli.assert_not_called()
        soffice.assert_called_once_with(pptx_path, out_dir, dpi=150)

    def test_officecli_failure_falls_back_to_soffice(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            out_dir = Path(directory) / "renders"
            with patch("scripts.presentation_qa.render_page._render_with_officecli", side_effect=RuntimeError("browser unavailable")) as officecli:
                with patch("scripts.presentation_qa.render_page._render_with_soffice", return_value=[out_dir / "slide-1.jpg"]) as soffice:
                    result = render_to_png(pptx_path, out_dir)

        self.assertEqual([out_dir / "slide-1.jpg"], result)
        officecli.assert_called_once_with(pptx_path, out_dir, dpi=150)
        soffice.assert_called_once_with(pptx_path, out_dir, dpi=150)

    def test_strict_officecli_renderer_does_not_fallback(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            out_dir = Path(directory) / "renders"
            with patch("scripts.presentation_qa.render_page._render_with_officecli", side_effect=RuntimeError("missing")):
                with patch("scripts.presentation_qa.render_page._render_with_soffice") as soffice:
                    with self.assertRaisesRegex(RuntimeError, "OfficeCLI render failed"):
                        render_to_png(pptx_path, out_dir, strict_renderer=True)

        soffice.assert_not_called()

    def test_officecli_renderer_uses_repository_font_browser_capture(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            out_dir = Path(directory) / "renders"
            self._write_text_deck(pptx_path)

            def run(command, **kwargs):
                if "html" in command:
                    return type(
                        "Completed",
                        (),
                        {"stdout": '<div class="slide">中文</div>', "stderr": "", "returncode": 0},
                    )()
                if command[0] == "/usr/bin/node":
                    screenshot = Path(command[-1])
                    Image.new("RGB", (960, 720), "white").save(screenshot, format="PNG")
                return type("Completed", (), {"stdout": "", "stderr": "", "returncode": 0})()

            with patch("scripts.presentation_qa.render_page._officecli_path", return_value=Path("/usr/bin/officecli")):
                with patch("scripts.presentation_qa.render_page.shutil.which", return_value="/usr/bin/node"):
                    with patch("scripts.presentation_qa.render_page.subprocess.run", side_effect=run) as calls:
                        result = _render_with_officecli(pptx_path, out_dir, dpi=150)

        self.assertEqual([out_dir / "slide-1.png"], result)
        commands = [call.args[0] for call in calls.call_args_list]
        html = next(command for command in commands if "html" in command)
        font_capture = next(command for command in commands if command[0] == "/usr/bin/node")
        self.assertEqual("/usr/bin/officecli", html[0])
        self.assertNotIn("stats", html)
        self.assertNotIn("-o", html)
        self.assertEqual(["--start", "1", "--end", "1"], html[-4:])
        self.assertTrue(font_capture[1].endswith("officecli_html_screenshot.mjs"))

    def test_fragmented_native_text_matches_one_script_string(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            self._write_text_deck(pptx_path)
            report = build_text_content_qa(
                pptx_path,
                ["生产运行需求：新能源接入后，安全保供需要及时数据支撑"],
                order_sensitive=False,
                allow_fragmented_actual=True,
            )

        self.assertTrue(report["valid"])
        self.assertEqual(2, len(report["actual_texts"]))

    def test_fragmented_native_text_ignores_space_lost_at_shape_boundary(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            for index, value in enumerate(("企业客户：", "单班，交易与", "AI 决策实训")):
                shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3 + index * 0.4), Inches(3), Inches(0.4))
                shape.text_frame.text = value
            presentation.save(str(pptx_path))

            report = build_text_content_qa(
                pptx_path,
                ["企业客户：单班，交易与 AI 决策实训"],
                order_sensitive=False,
                allow_fragmented_actual=True,
            )

        self.assertTrue(report["valid"])


if __name__ == "__main__":
    unittest.main()
