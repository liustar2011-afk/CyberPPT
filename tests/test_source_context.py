from __future__ import annotations

import json
from pathlib import Path
import sys
from types import ModuleType, SimpleNamespace

import pytest
from pptx import Presentation
from pptx.util import Inches

from cyberppt.source_document_map import (
    SCRIPT_SOURCE_INDEX,
    load_source_units,
    prepare_source_context,
    prepare_source_map,
)
from script_engine.source_index import estimate_reading_load, recommend_reading_mode


def test_script_source_context_is_deterministic_and_writes_only_one_cache(tmp_path: Path) -> None:
    project = tmp_path / "project"
    source = project / "source"
    source.mkdir(parents=True)
    (source / "brief.md").write_text("# 总体判断\n关键证据。\n", encoding="utf-8")

    first = prepare_source_context(project)
    first_bytes = (project / SCRIPT_SOURCE_INDEX).read_bytes()
    second = prepare_source_context(project)

    assert first["schema"] == "cyberppt.source_index.v2"
    assert first["status"] == "passed"
    assert first["source_hashes"] == second["source_hashes"]
    assert first_bytes == (project / SCRIPT_SOURCE_INDEX).read_bytes()
    assert not (project / "workbench/stages/00-source-map").exists()
    generated = [
        path.relative_to(project).as_posix()
        for path in project.rglob("*")
        if path.is_file() and not path.is_relative_to(source)
    ]
    assert generated == ["script/.cache/source-index.json"]


def test_multi_file_reading_load_enters_long_mode_without_dropping_units(tmp_path: Path) -> None:
    project = tmp_path / "project"
    source = project / "sources"
    source.mkdir(parents=True)
    (source / "a.txt").write_text("甲" * 31_000, encoding="utf-8")
    (source / "b.txt").write_text("乙" * 31_000, encoding="utf-8")

    payload = prepare_source_context(project)

    assert payload["reading_recommendation"]["mode"] == "long"
    assert payload["reading_load"]["source_count"] == 2
    assert sum(len(item["text"]) for item in payload["units"]) == 62_000
    assert payload["reading_strategy"]["excluded_unit_ids"] == []


def test_reading_mode_boundary_keeps_45_pages_direct_and_routes_46_to_long() -> None:
    def units(count: int) -> list[dict]:
        return [
            {
                "unit_id": f"SU-{index}",
                "source_id": "SRC-1",
                "text": "简短内容",
                "locator": {"slide": index},
            }
            for index in range(1, count + 1)
        ]

    direct = recommend_reading_mode(estimate_reading_load(units(45), [{"source_id": "SRC-1"}]))
    long = recommend_reading_mode(estimate_reading_load(units(46), [{"source_id": "SRC-1"}]))

    assert direct["mode"] == "direct"
    assert long["mode"] == "long"


def test_pptx_native_extractor_keeps_title_text_table_and_notes(tmp_path: Path) -> None:
    project = tmp_path / "project"
    source = project / "source"
    source.mkdir(parents=True)
    path = source / "brief.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "总体方案"
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "核心证据"
    table = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "指标"
    table.cell(0, 1).text = "数值"
    slide.notes_slide.notes_text_frame.text = "演讲备注"
    deck.save(path)

    payload = prepare_source_context(project)
    kinds = {item["kind"] for item in payload["units"]}
    text = "\n".join(item["text"] for item in payload["units"])

    assert payload["status"] == "passed"
    assert {"heading", "paragraph", "table_row", "speaker_note"}.issubset(kinds)
    assert "总体方案" in text
    assert "核心证据" in text
    assert "指标 | 数值" in text
    assert "演讲备注" in text

    strict = prepare_source_map(project)
    assert strict["status"] == "passed"
    assert [item["unit_id"] for item in payload["units"]] == [
        item["unit_id"] for item in load_source_units(project)
    ]


def test_xlsx_native_extractor_is_used_when_openpyxl_is_available(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    project = tmp_path / "project"
    source = project / "source"
    source.mkdir(parents=True)
    path = source / "metrics.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "指标表"
    sheet.append(["指标", "目标"])
    sheet.append(["覆盖率", "95%"])
    workbook.save(path)

    payload = prepare_source_context(project)
    text = "\n".join(item["text"] for item in payload["units"])

    assert payload["status"] == "passed"
    assert "指标表" in text
    assert "覆盖率 | 95%" in text
    assert not any(item["kind"] == "binary" for item in payload["units"])


def test_pdf_uses_format_level_markitdown_fallback_when_available(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    fake_module = ModuleType("markitdown")

    class FakeMarkItDown:
        def convert(self, _path: str) -> SimpleNamespace:
            return SimpleNamespace(text_content="# 政策依据\n精确条款内容。")

    fake_module.MarkItDown = FakeMarkItDown  # type: ignore[attr-defined]
    monkeypatch.setitem(sys.modules, "markitdown", fake_module)
    project = tmp_path / "project"
    source = project / "source"
    source.mkdir(parents=True)
    (source / "policy.pdf").write_bytes(b"%PDF-fallback-fixture")

    payload = prepare_source_context(project)

    assert payload["status"] == "passed"
    assert "政策依据" in "\n".join(item["text"] for item in payload["units"])
    assert any(
        item["code"] == "SOURCE_FORMAT_EXTRACTED_WITH_MARKITDOWN"
        for item in payload["warnings"]
    )


def test_source_index_is_json_round_trippable(tmp_path: Path) -> None:
    project = tmp_path / "project"
    source = project / "source"
    source.mkdir(parents=True)
    (source / "brief.txt").write_text("可追溯内容", encoding="utf-8")

    prepare_source_context(project)
    payload = json.loads((project / SCRIPT_SOURCE_INDEX).read_text(encoding="utf-8"))

    assert payload["schema"] == "cyberppt.source_index.v2"
    assert payload["units"][0]["unit_id"].startswith("SU-")
