from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from cyberppt.commands.production_qa import render_and_compare, validate_assembly_bundle
from scripts.validate_pptx import validate_pptx


def _write_json(path: Path, payload: dict) -> Path:
    path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    return path


def _write_minimal_pptx(path: Path) -> Path:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr(
            "ppt/presentation.xml",
            '<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:sldSz cx="12192000" cy="6858000"/></p:presentation>',
        )
        archive.writestr("ppt/slides/slide1.xml", "<p:sld/>")
        archive.writestr("ppt/notesSlides/notesSlide1.xml", "<p:notes/>")
    return path


def _write_full_image_pptx(path: Path, image_path: Path, *, notes: bool = True, title: bool = True) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    if title:
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4.0), Inches(0.4)).text = "Native title"
    slide.shapes.add_picture(str(image_path), Inches(0), Inches(0.8), width=Inches(13.333), height=Inches(6.2))
    presentation.save(path)
    if notes:
        with zipfile.ZipFile(path, "a") as archive:
            archive.writestr("ppt/notesSlides/notesSlide1.xml", "<p:notes/>")
    return path


class ProductionQaTests(unittest.TestCase):
    def test_validates_complete_approved_assembly_bundle(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            full = (root / "page_001_full.png")
            full.write_bytes(b"approved-image")
            template_manifest = _write_json(
                root / "template_image_manifest.json",
                {"tasks": [{"page_number": 1, "image_path": str(full), "notes_text": "approved notes"}]},
            )
            pptx = _write_minimal_pptx(root / "assembled.pptx")

            report = validate_assembly_bundle(
                {
                    "project": str(root),
                    "exported_pptx": str(pptx),
                    "template_image_manifest": str(template_manifest),
                    "approved_images": {1: str(full)},
                },
                [1],
            )

        self.assertTrue(report["valid"])
        self.assertTrue(report["checks"]["pptx_readable"])
        self.assertTrue(report["checks"]["notes_complete"])

    def test_rejects_assembly_bundle_outside_project(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory) / "project"
            root.mkdir()
            outside = Path(directory) / "outside.pptx"
            _write_minimal_pptx(outside)

            report = validate_assembly_bundle(
                {
                    "project": str(root),
                    "exported_pptx": str(outside),
                    "template_image_manifest": str(outside),
                    "approved_images": {},
                },
                [1],
            )

        self.assertFalse(report["valid"])
        self.assertIn("output_outside_project", report["failures"])

    def test_render_and_compare_blocks_when_render_tool_unavailable(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            pptx = _write_minimal_pptx(root / "assembled.pptx")
            manifest = _write_json(
                root / "template_image_manifest.json",
                {
                    "canvas": {"width": 1280, "height": 720},
                    "body_region": {"x": 0, "y": 0, "width": 1280, "height": 720},
                    "tasks": [{"page_number": 1, "image_path": str(root / "full.png"), "notes_text": "notes"}],
                },
            )

            with patch("cyberppt.commands.production_qa.render_to_png", return_value=[]):
                with self.assertRaisesRegex(RuntimeError, "render_tool_unavailable"):
                    render_and_compare(pptx, manifest, {1: root / "full.png"}, [1], root / "qa")

    def test_render_and_compare_passes_matching_body_image(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            approved = root / "full.png"
            rendered = root / "rendered.jpg"
            Image.new("RGB", (40, 20), "#336699").save(approved)
            Image.new("RGB", (40, 20), "#336699").save(rendered)
            pptx = _write_minimal_pptx(root / "assembled.pptx")
            manifest = _write_json(
                root / "template_image_manifest.json",
                {
                    "canvas": {"width": 1280, "height": 720},
                    "body_region": {"x": 0, "y": 0, "width": 1280, "height": 720},
                    "tasks": [{"page_number": 1, "image_path": str(approved), "notes_text": "notes"}],
                },
            )

            with patch("cyberppt.commands.production_qa.render_to_png", return_value=[rendered]):
                report = render_and_compare(pptx, manifest, {1: approved}, [1], root / "qa")

        self.assertTrue(report["passed"])
        self.assertLessEqual(report["slides"][0]["mean_abs_diff"], 1.0)

    def test_render_and_compare_rejects_large_body_difference(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            approved = root / "full.png"
            rendered = root / "rendered.jpg"
            Image.new("RGB", (40, 20), "#000000").save(approved)
            Image.new("RGB", (40, 20), "#ffffff").save(rendered)
            pptx = _write_minimal_pptx(root / "assembled.pptx")
            manifest = _write_json(
                root / "template_image_manifest.json",
                {
                    "canvas": {"width": 1280, "height": 720},
                    "body_region": {"x": 0, "y": 0, "width": 1280, "height": 720},
                    "tasks": [{"page_number": 1, "image_path": str(approved), "notes_text": "notes"}],
                },
            )

            with patch("cyberppt.commands.production_qa.render_to_png", return_value=[rendered]):
                report = render_and_compare(pptx, manifest, {1: approved}, [1], root / "qa")

        self.assertFalse(report["passed"])
        self.assertGreater(report["slides"][0]["mean_abs_diff"], report["slides"][0]["threshold"])

    def test_strict_validator_accepts_full_image_delivery_manifest(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            approved = root / "approved.png"
            Image.new("RGB", (100, 60), "#336699").save(approved)
            pptx = _write_full_image_pptx(root / "deck.pptx", approved)
            visual = _write_json(root / "production_visual_report.json", {"passed": True})
            lock = _write_json(
                root / "template_text_lock.json",
                {"records": [{"page": 1, "title": "Native title", "approved": True}]},
            )
            manifest = _write_json(
                root / "full_image_delivery_manifest.json",
                {
                    "schema": "cyberppt.full_image_delivery_manifest.v1",
                    "delivery_mode": "full_image_ppt",
                    "body_content_editable": False,
                    "template_text_editable": True,
                    "speaker_notes_required": True,
                    "template_text_lock": {"path": str(lock), "sha256": _sha256_for_test(lock)},
                    "production_visual_report": {"path": str(visual), "passed": True},
                    "slides": [
                        {
                            "slide": 1,
                            "delivery_mode": "full_image_ppt",
                            "native_text_requirements": ["Native title"],
                            "image_assets": [{"role": "approved_full_image", "path": str(approved)}],
                        }
                    ],
                },
            )

            report = validate_pptx(pptx, manifest_path=manifest, strict=True)

        self.assertEqual([], report["errors"])

    def test_strict_validator_rejects_full_image_without_notes(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            approved = root / "approved.png"
            Image.new("RGB", (100, 60), "#336699").save(approved)
            pptx = _write_full_image_pptx(root / "deck.pptx", approved, notes=False)
            visual = _write_json(root / "production_visual_report.json", {"passed": True})
            lock = _write_json(
                root / "template_text_lock.json",
                {"records": [{"page": 1, "title": "Native title", "approved": True}]},
            )
            manifest = _write_json(
                root / "full_image_delivery_manifest.json",
                {
                    "schema": "cyberppt.full_image_delivery_manifest.v1",
                    "delivery_mode": "full_image_ppt",
                    "body_content_editable": False,
                    "template_text_editable": True,
                    "speaker_notes_required": True,
                    "template_text_lock": {"path": str(lock), "sha256": _sha256_for_test(lock)},
                    "production_visual_report": {"path": str(visual), "passed": True},
                    "slides": [
                        {
                            "slide": 1,
                            "delivery_mode": "full_image_ppt",
                            "native_text_requirements": ["Native title"],
                            "image_assets": [{"role": "approved_full_image", "path": str(approved)}],
                        }
                    ],
                },
            )

            report = validate_pptx(pptx, manifest_path=manifest, strict=True)

        self.assertIn("FULL_IMAGE_SPEAKER_NOTES_MISSING", {item["code"] for item in report["errors"]})

    def test_strict_validator_rejects_full_image_with_wrong_template_title(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            approved = root / "approved.png"
            Image.new("RGB", (100, 60), "#336699").save(approved)
            pptx = _write_full_image_pptx(root / "deck.pptx", approved)
            visual = _write_json(root / "production_visual_report.json", {"passed": True})
            lock = _write_json(
                root / "template_text_lock.json",
                {"records": [{"page": 1, "title": "Approved title", "approved": True}]},
            )
            manifest = _write_json(
                root / "full_image_delivery_manifest.json",
                {
                    "schema": "cyberppt.full_image_delivery_manifest.v1",
                    "delivery_mode": "full_image_ppt",
                    "body_content_editable": False,
                    "template_text_editable": True,
                    "speaker_notes_required": True,
                    "template_text_lock": {"path": str(lock), "sha256": _sha256_for_test(lock)},
                    "production_visual_report": {"path": str(visual), "passed": True},
                    "slides": [
                        {
                            "slide": 1,
                            "delivery_mode": "full_image_ppt",
                            "native_text_requirements": ["Approved title"],
                            "image_assets": [{"role": "approved_full_image", "path": str(approved)}],
                        }
                    ],
                },
            )

            report = validate_pptx(pptx, manifest_path=manifest, strict=True)

        self.assertIn("FULL_IMAGE_NATIVE_TEMPLATE_TEXT_MISSING", {item["code"] for item in report["errors"]})


def _sha256_for_test(path: Path) -> str:
    import hashlib

    return hashlib.sha256(path.read_bytes()).hexdigest()


if __name__ == "__main__":
    unittest.main()
